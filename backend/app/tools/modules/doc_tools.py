import os
import json
import yaml
from typing import Dict, Any, List
from app.tools.base import BaseTool
from app.logging.logger import get_logger

logger = get_logger(__name__)

WORKSPACE_DIR = os.path.abspath(os.path.join(os.getcwd(), "workspace"))
os.makedirs(WORKSPACE_DIR, exist_ok=True)

class DocxTool(BaseTool):
    @property
    def slug(self) -> str: return "docx_tool"
    @property
    def name(self) -> str: return "Word Document Processor (.docx)"
    @property
    def description(self) -> str:
        return "Reads, creates, or appends to Microsoft Word .docx files in workspace. Arguments: action ('read'/'create'/'append'), filename, content (text or paragraph content)."

    async def run(self, **kwargs) -> str:
        action = str(kwargs.get("action", "")).strip().lower()
        filename = kwargs.get("filename") or kwargs.get("file") or kwargs.get("name") or ""
        filename = str(filename).strip()
        content = kwargs.get("content") or kwargs.get("text") or kwargs.get("data") or ""
        content = str(content).strip()

        if action not in ["read", "create", "append"]:
            return "Error: action must be 'read', 'create', or 'append'."
        if not filename:
            return "Error: filename parameter is required."
        if not filename.endswith(".docx"):
            filename += ".docx"

        filepath = os.path.join(WORKSPACE_DIR, os.path.basename(filename))

        try:
            import docx
            if action == "read":
                if not os.path.exists(filepath):
                    return f"Error: File '{filename}' not found in workspace."
                doc = docx.Document(filepath)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                tables_data = []
                for t in doc.tables:
                    for r in t.rows:
                        tables_data.append(" | ".join([cell.text.strip() for cell in r.cells]))
                res = f"Contents of {filename}:\n\n" + "\n".join(paragraphs)
                if tables_data:
                    res += "\n\nTables:\n" + "\n".join(tables_data)
                return res

            elif action == "create":
                doc = docx.Document()
                doc.add_heading(os.path.splitext(os.path.basename(filename))[0].replace("_", " ").title(), level=1)
                for line in content.split("\n"):
                    if line.startswith("# "):
                        doc.add_heading(line[2:], level=1)
                    elif line.startswith("## "):
                        doc.add_heading(line[3:], level=2)
                    elif line.startswith("- "):
                        doc.add_paragraph(line[2:], style='List Bullet')
                    elif line.strip():
                        doc.add_paragraph(line)
                doc.save(filepath)
                return f"Successfully created Word document '{filename}' in workspace."

            elif action == "append":
                if os.path.exists(filepath):
                    doc = docx.Document(filepath)
                else:
                    doc = docx.Document()
                for line in content.split("\n"):
                    if line.strip():
                        doc.add_paragraph(line)
                doc.save(filepath)
                return f"Successfully appended content to '{filename}'."

        except Exception as e:
            return f"DocxTool error: {str(e)}"

class ExcelWriterTool(BaseTool):
    @property
    def slug(self) -> str: return "excel_writer_tool"
    @property
    def name(self) -> str: return "Excel Spreadsheet Builder (.xlsx)"
    @property
    def description(self) -> str:
        return "Creates or updates Excel spreadsheets (.xlsx) in workspace. Arguments: filename, data (JSON list of dicts/lists or CSV formatted string), sheet_name (optional, default 'Sheet1')."

    async def run(self, **kwargs) -> str:
        filename = kwargs.get("filename") or kwargs.get("file") or kwargs.get("name") or ""
        filename = str(filename).strip()
        data_input = kwargs.get("data") or kwargs.get("content") or kwargs.get("rows") or ""
        sheet_name = str(kwargs.get("sheet_name", "Sheet1")).strip()

        if not filename:
            return "Error: filename parameter is required."
        if not filename.endswith(".xlsx"):
            filename += ".xlsx"

        filepath = os.path.join(WORKSPACE_DIR, os.path.basename(filename))

        try:
            import pandas as pd
            parsed_data = None
            if isinstance(data_input, list):
                parsed_data = data_input
            elif isinstance(data_input, str):
                data_str = data_input.strip()
                if data_str.startswith("[") or data_str.startswith("{"):
                    parsed_data = json.loads(data_str)
                else:
                    # CSV string format fallback
                    lines = [l.split(",") for l in data_str.splitlines() if l.strip()]
                    if lines:
                        headers = [h.strip() for h in lines[0]]
                        rows = []
                        for row in lines[1:]:
                            rows.append({headers[i]: row[i].strip() if i < len(row) else "" for i in range(len(headers))})
                        parsed_data = rows

            if not parsed_data:
                return "Error: Could not parse data input. Provide a valid JSON array of objects or CSV text."

            df = pd.DataFrame(parsed_data)
            
            # Save using openpyxl
            with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                df.to_excel(writer, sheet_name=sheet_name, index=False)

            return f"Successfully created Excel spreadsheet '{filename}' ({len(df)} rows, {len(df.columns)} columns) in workspace."

        except Exception as e:
            return f"ExcelWriterTool error: {str(e)}"

class PptxTool(BaseTool):
    @property
    def slug(self) -> str: return "pptx_tool"
    @property
    def name(self) -> str: return "PowerPoint Presentation Generator (.pptx)"
    @property
    def description(self) -> str:
        return "Reads or creates PowerPoint presentation decks (.pptx) in workspace. Arguments: action ('read'/'create'), filename, slides (JSON array of slide dicts with 'title' and 'bullets'/'content' for create)."

    async def run(self, **kwargs) -> str:
        action = str(kwargs.get("action", "")).strip().lower()
        filename = kwargs.get("filename") or kwargs.get("file") or kwargs.get("name") or ""
        filename = str(filename).strip()
        slides_input = kwargs.get("slides") or kwargs.get("content") or kwargs.get("data") or ""

        if action not in ["read", "create"]:
            return "Error: action must be 'read' or 'create'."
        if not filename:
            return "Error: filename parameter is required."
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        filepath = os.path.join(WORKSPACE_DIR, os.path.basename(filename))

        try:
            import pptx
            if action == "read":
                if not os.path.exists(filepath):
                    return f"Error: File '{filename}' not found in workspace."
                prs = pptx.Presentation(filepath)
                slides_out = []
                for idx, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text.strip())
                    slides_out.append(f"Slide {idx}:\n" + "\n".join(texts))
                return f"PowerPoint Deck ({filename}) - {len(prs.slides)} slides:\n\n" + "\n\n---\n\n".join(slides_out)

            elif action == "create":
                prs = pptx.Presentation()
                slides_list = []
                if isinstance(slides_input, list):
                    slides_list = slides_input
                elif isinstance(slides_input, str) and slides_input.strip():
                    try:
                        slides_list = json.loads(slides_input)
                    except Exception:
                        # Simple split fallback
                        slides_list = [{"title": "Presentation", "bullets": slides_input.splitlines()}]

                if not slides_list:
                    slides_list = [{"title": "AgentHive Presentation", "bullets": ["Generated via AgentHive PPTX Tool"]}]

                title_slide_layout = prs.slide_layouts[0]
                bullet_slide_layout = prs.slide_layouts[1]

                for idx, s in enumerate(slides_list):
                    layout = title_slide_layout if idx == 0 else bullet_slide_layout
                    slide = prs.slides.add_slide(layout)
                    
                    title_box = slide.shapes.title
                    if title_box:
                        title_box.text = s.get("title", f"Slide {idx + 1}")

                    if idx > 0 and len(slide.shapes) > 1:
                        tf = slide.shapes[1].text_frame
                        bullets = s.get("bullets", s.get("content", []))
                        if isinstance(bullets, str):
                            bullets = bullets.splitlines()
                        for b_idx, bullet_text in enumerate(bullets):
                            if b_idx == 0:
                                tf.text = str(bullet_text)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(bullet_text)

                prs.save(filepath)
                return f"Successfully created PowerPoint presentation '{filename}' ({len(slides_list)} slides) in workspace."

        except Exception as e:
            return f"PptxTool error: {str(e)}"

class OcrTool(BaseTool):
    @property
    def slug(self) -> str: return "ocr_tool"
    @property
    def name(self) -> str: return "Image OCR Text Extractor"
    @property
    def description(self) -> str:
        return "Extracts text from image files (PNG, JPG, TIFF) in workspace using OCR with graceful fallback. Arguments: filename."

    async def run(self, **kwargs) -> str:
        filename = kwargs.get("filename") or kwargs.get("file") or kwargs.get("name") or kwargs.get("image") or ""
        filename = str(filename).strip()
        if not filename:
            return "Error: filename parameter is required."

        filepath = os.path.join(WORKSPACE_DIR, os.path.basename(filename))
        if not os.path.exists(filepath):
            return f"Error: Image '{filename}' not found in workspace."

        try:
            from PIL import Image
            img = Image.open(filepath)
            format_info = f"Format: {img.format}, Size: {img.width}x{img.height}, Mode: {img.mode}"

            try:
                import pytesseract
                text = pytesseract.image_to_string(img)
                if text.strip():
                    return f"OCR Text Extracted from '{filename}' ({format_info}):\n\n{text}"
                else:
                    return f"OCR processed '{filename}' ({format_info}), but no text was detected."
            except Exception as tess_err:
                return (
                    f"OCR Fallback for '{filename}':\n"
                    f"{format_info}\n"
                    f"Note: Full text OCR requires system Tesseract binary (error: {str(tess_err)})."
                )
        except Exception as e:
            return f"OcrTool error: {str(e)}"

class JsonYamlTool(BaseTool):
    @property
    def slug(self) -> str: return "json_yaml_tool"
    @property
    def name(self) -> str: return "JSON / YAML Converter & Validator"
    @property
    def description(self) -> str:
        return "Validates, formats, or converts between JSON and YAML. Arguments: action ('json_to_yaml'/'yaml_to_json'/'validate_json'/'validate_yaml'/'format_json'), content."

    async def run(self, **kwargs) -> str:
        action = str(kwargs.get("action", "")).strip().lower()
        content = kwargs.get("content") or kwargs.get("data") or kwargs.get("text") or kwargs.get("json") or kwargs.get("yaml") or ""
        content = str(content).strip()

        if not action or not content:
            return "Error: Both 'action' and 'content' parameters are required."

        try:
            if action == "json_to_yaml":
                parsed = json.loads(content)
                converted_yaml = yaml.dump(parsed, sort_keys=False, default_flow_style=False)
                return f"Converted YAML:\n```yaml\n{converted_yaml}\n```"

            elif action == "yaml_to_json":
                parsed = yaml.safe_load(content)
                converted_json = json.dumps(parsed, indent=2)
                return f"Converted JSON:\n```json\n{converted_json}\n```"

            elif action == "validate_json":
                json.loads(content)
                return "Valid JSON document."

            elif action == "validate_yaml":
                yaml.safe_load(content)
                return "Valid YAML document."

            elif action == "format_json":
                parsed = json.loads(content)
                formatted = json.dumps(parsed, indent=2)
                return f"Formatted JSON:\n```json\n{formatted}\n```"

            else:
                return f"Error: Invalid action '{action}'. Supported: json_to_yaml, yaml_to_json, validate_json, validate_yaml, format_json."
        except Exception as e:
            return f"JsonYamlTool error: Invalid format or parsing error: {str(e)}"
